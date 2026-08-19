"""Estilo visual de las presentaciones del curso INF-396.

Reproduce el diseño de la clase 01, para que todas las clases se vean iguales.
Los scripts build_claseNN_ppt.py importan estas funciones.
"""

from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt

# Paleta del curso
AZUL_PORTADA = RGBColor(0x3B, 0x4B, 0x9E)
AZUL_SECCION = RGBColor(0x1D, 0x35, 0x57)
DORADO = RGBColor(0xE0, 0xA4, 0x37)
FONDO = RGBColor(0xFA, 0xFB, 0xFC)
TITULO = RGBColor(0x0F, 0x0F, 0x0F)
CUERPO = RGBColor(0x3A, 0x4A, 0x5C)
SECUNDARIO = RGBColor(0x6B, 0x77, 0x85)
BLANCO = RGBColor(0xFF, 0xFF, 0xFF)
CLARO = RGBColor(0xB8, 0xC4, 0xD4)
CODIGO = RGBColor(0x1D, 0x35, 0x57)

FUENTE = "PT Serif"
FUENTE_CODIGO = "Menlo"

ANCHO = Inches(10)
ALTO = Inches(5.62)


def _fondo(slide, color):
    forma = slide.shapes.add_shape(1, 0, 0, ANCHO, ALTO)
    forma.fill.solid()
    forma.fill.fore_color.rgb = color
    forma.line.fill.background()
    forma.shadow.inherit = False
    return forma


def _caja(slide, izq, arriba, ancho, alto):
    caja = slide.shapes.add_textbox(
        Inches(izq), Inches(arriba), Inches(ancho), Inches(alto)
    )
    caja.text_frame.word_wrap = True
    return caja


def _run(parrafo, texto, tam, color, negrita=False, fuente=FUENTE):
    run = parrafo.add_run()
    run.text = texto
    run.font.name = fuente
    run.font.size = Pt(tam)
    run.font.bold = negrita
    run.font.color.rgb = color
    return run


def portada(prs, titulo, subtitulo, autor, meta):
    """Primera slide: fondo azul, título grande y chip dorado con el número de clase."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _fondo(slide, AZUL_PORTADA)

    caja = _caja(slide, 0.77, 1.37, 8.47, 1.42)
    _run(caja.text_frame.paragraphs[0], titulo, 42, BLANCO, negrita=True)

    chip = slide.shapes.add_shape(1, Inches(2.16), Inches(2.9), Inches(5.69), Inches(0.57))
    chip.fill.solid()
    chip.fill.fore_color.rgb = DORADO
    chip.line.fill.background()
    chip.shadow.inherit = False
    _run(chip.text_frame.paragraphs[0], subtitulo, 19, AZUL_SECCION, negrita=True)

    caja = _caja(slide, 0.77, 3.88, 8.47, 0.98)
    tf = caja.text_frame
    _run(tf.paragraphs[0], autor, 15, BLANCO)
    _run(tf.add_paragraph(), meta, 13, CLARO)
    return slide


def seccion(prs, titulo, bajada):
    """Slide divisoria entre partes de la clase."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _fondo(slide, AZUL_SECCION)
    caja = _caja(slide, 0.66, 2.02, 8.69, 1.53)
    tf = caja.text_frame
    _run(tf.paragraphs[0], titulo, 34, BLANCO, negrita=True)
    _run(tf.add_paragraph(), bajada, 18, DORADO)
    return slide


def _base_contenido(prs, titulo):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _fondo(slide, FONDO)
    caja = _caja(slide, 0.44, 0.33, 9.13, 0.87)
    _run(caja.text_frame.paragraphs[0], titulo, 28, TITULO, negrita=True)
    return slide


def _escribir_puntos(tf, puntos, tam, tam_sub):
    """Escribe la lista de puntos. Un punto es texto, o una tupla (texto, destacado),
    o una lista de sub puntos."""
    primero = True
    for punto in puntos:
        if isinstance(punto, list):
            for sub in punto:
                p = tf.paragraphs[0] if primero else tf.add_paragraph()
                primero = False
                _run(p, "     ◦  ", tam_sub, SECUNDARIO)
                _run(p, sub, tam_sub, SECUNDARIO)
            continue

        texto, destacado = punto if isinstance(punto, tuple) else (punto, False)
        p = tf.paragraphs[0] if primero else tf.add_paragraph()
        primero = False
        _run(p, "▸ ", tam, DORADO, negrita=True)
        _run(p, texto, tam, AZUL_SECCION if destacado else CUERPO, negrita=destacado)


def contenido(prs, titulo, puntos):
    """Slide de texto a ancho completo."""
    slide = _base_contenido(prs, titulo)
    caja = _caja(slide, 0.77, 1.48, 8.47, 3.61)
    _escribir_puntos(caja.text_frame, puntos, 18, 14)
    return slide


def contenido_con_figura(prs, titulo, puntos, imagen):
    """Slide con texto a la izquierda y una figura a la derecha."""
    slide = _base_contenido(prs, titulo)
    caja = _caja(slide, 0.55, 1.37, 4.4, 3.5)
    _escribir_puntos(caja.text_frame, puntos, 15, 13)
    slide.shapes.add_picture(str(imagen), Inches(5.35), Inches(1.5), width=Inches(4.2))
    return slide


def figura_grande(prs, titulo, imagen, pie=None):
    """Slide dedicada a una figura, con un pie de foto opcional."""
    slide = _base_contenido(prs, titulo)
    slide.shapes.add_picture(str(imagen), Inches(1.55), Inches(1.3), width=Inches(6.9))
    if pie:
        caja = _caja(slide, 0.77, 4.85, 8.47, 0.5)
        _run(caja.text_frame.paragraphs[0], pie, 13, SECUNDARIO)
    return slide


def codigo(prs, titulo, lineas, comentario=None):
    """Slide con un bloque de código corto."""
    slide = _base_contenido(prs, titulo)

    panel = slide.shapes.add_shape(1, Inches(0.77), Inches(1.4), Inches(8.47), Inches(0.55 * len(lineas) + 0.4))
    panel.fill.solid()
    panel.fill.fore_color.rgb = RGBColor(0xF0, 0xF2, 0xF5)
    panel.line.color.rgb = RGBColor(0xD8, 0xDD, 0xE3)
    panel.shadow.inherit = False
    panel.text_frame.word_wrap = True

    tf = panel.text_frame
    tf.margin_left = Inches(0.25)
    tf.margin_top = Inches(0.15)
    for i, linea in enumerate(lineas):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = 1 if False else None
        _run(p, linea, 15, CODIGO, fuente=FUENTE_CODIGO)

    if comentario:
        caja = _caja(slide, 0.77, 1.4 + 0.55 * len(lineas) + 0.65, 8.47, 1.2)
        _escribir_puntos(caja.text_frame, [comentario], 16, 13)
    return slide
